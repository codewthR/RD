import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from .SlideFormating import add_bullets,add_image_placeholder,add_numbered_shapes,add_paragraph,add_textbox,adding_title,rgb_from_string


# ========================================================= Slide Creation ===================================================

from pptx.enum.text import PP_ALIGN

def create_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Blank layout

    # Presentation Topic
    if slide_data.get("presentation_topic"):
        text = slide_data["presentation_topic"]
        color = rgb_from_string(slide_data.get("presentation_topic_color_value", "(0,0,0)"))
        font_size = int(slide_data.get("presentation_topic_font_size", 36))
        box = slide_data["Inches"].get("presentation_topic_inches", [0, 0, 0, 0])
        add_textbox(slide, *box, text=text, font_size=font_size, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Presentation Subtopic
    if slide_data.get("presentation_subtopic"):
        text = slide_data["presentation_subtopic"]
        color = rgb_from_string(slide_data.get("presentation_topic_color_value", "(0,0,0)"))
        font_size = int(slide_data.get("text_font_value", 24))
        box = slide_data["Inches"].get("presentation_subtopic", [0, 0, 0, 0])
        add_textbox(slide, *box, text=text, font_size=font_size, bold=False, color=color, align=PP_ALIGN.CENTER)

    # Slide Title
    if slide_data.get("title"):
        text = slide_data["title"]
        font_size = int(slide_data.get("title_font_value", 28))
        color = rgb_from_string(slide_data.get("title_color_value", "(0,0,0)"))
        box = slide_data["Inches"].get("title_Inches", [0, 0, 0, 0])
        add_textbox(slide, *box, text=text, font_size=font_size, bold=True, color=color)

    # Slide Content
    content_type = slide_data.get("contentType", "")
    content = slide_data.get("content", {})
    content_box = slide_data["Inches"].get("content_Inches", [0, 0, 0, 0])
    text_font_color = rgb_from_string(slide_data.get("all_text_color_", "(0,0,0)"))
    text_font_size = int(slide_data.get("text_font_value", 20))

    if content_type == "paragraph":
        add_paragraph(slide, content.get("paragraph", ""), content_box, font_size=text_font_size, color=text_font_color)

    elif content_type == "bullets":
        bullet_box = slide_data["Inches"].get("bullets_shapes_inches", [0, 0, 0, 0])
        add_bullets(slide, content.get("bullets", []),bullet_box, font_size=text_font_size, color=text_font_color)

    elif content_type == "numbered":
        numbered_box = slide_data["Inches"].get("numbered_shapes_inches", [0, 0, 0, 0])
        add_numbered_shapes(slide, content.get("numbered", []), numbered_box, font_size=text_font_size, color=text_font_color)

    elif content_type == "numbered":
        numbered_box = slide_data["Inches"].get("numbered_shapes_inches", [0, 0, 0, 0])
        add_numbered_shapes(slide, content.get("numbered", []), numbered_box, font_size=text_font_size, color=text_font_color)

    # Slide Image
    if slide_data.get("imagePrompt"):
        image_box = slide_data["Inches"].get("Images_Inches", [0, 0, 0, 0])
        add_image_placeholder(slide, slide_data["imagePrompt"], image_box)


