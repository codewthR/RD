import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

from .Text_cleaning import basic_clean

# left, top , width ,height

# left Distance from the left edge of the slide to the left edge of the textbox.
# top	Inches(0.3)	Distance from the top edge of the slide to the top edge of the textbox.
# width	Inches(9)	Width of the textbox.
# height	Inches(1)	Height of the textbox.

# colors tuple in fromat of (122,23,234)

# Function to add title text box on a given slide
def adding_title(slide, title, left, top, width, height, font_size, title_colors_tuple):
    title_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.vertical_anchor = MSO_ANCHOR.TOP
    title_frame.text = title
    p = title_frame.paragraphs[0]
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*title_colors_tuple)


def paragraph_add(slide, paragraph, left, top, width, height, para_font_size, para_colors_tuple):
    # Creating a textbox shape on the slide
    para_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    para_frame = para_box.text_frame
    para_frame.clear()  # Clear any existing content
    para_frame.word_wrap = True  # Wrap text within the box
    para_frame.vertical_anchor = MSO_ANCHOR.TOP  # Align text to top of box
    # Add paragraph and set text
    p = para_frame.add_paragraph()
    p.text = paragraph
    p.font.size = Pt(para_font_size)
    p.font.color.rgb = RGBColor(*para_colors_tuple)  # Unpack RGB values
    para_frame.margin_left = Inches(0.1)
    para_frame.margin_right = Inches(0.1)



def bullets_adding(slide, bullets, left, top, width, height, bullet_font_size, bullets_colors_tuple):
    bullet_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    bullet_frame = bullet_box.text_frame
    bullet_frame.clear()
    bullet_frame.word_wrap = True
    bullet_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Add each bullet as a paragraph 
    for bullet in bullets:
        cleaned_bullet = basic_clean(bullet)  # Basic clean is used to removing outliers from the data 
        p = bullet_frame.add_paragraph()
        p.text = "➕ " + cleaned_bullet
        p.level = 0
        p.font.size = Pt(bullet_font_size)
        p.font.color.rgb = RGBColor(*bullets_colors_tuple)  # Unpack tuple for the use case of giving colors



def shapes_adding(slide, y_offset, full_text):
    # Add a rounded rectangle shape to the slide
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.0), Inches(y_offset), Inches(5.5), Inches(0.8)
    )

    # Set shape fill color (light lavender)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 230, 250)

    # Set shape border color
    line = shape.line
    line.color.rgb = RGBColor(180, 180, 180)

    # Add and format text inside the shape
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = full_text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(60, 60, 60)

    return shape







# ============================================================================================================
# ====================================================== Utility Functions ===================================
# ============================================================================================================

def rgb_from_string(rgb_str):
    try:
        r, g, b = map(int, rgb_str.strip("() ").split(","))
        return RGBColor(r, g, b)
    except:
        return RGBColor(0, 0, 0)

def add_textbox(slide, x, y, w, h, text, font_size=20, bold=False, color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
    textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)-1)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    return textbox




from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def add_bullets(slide, items,bullet_box, font_size=18, color=RGBColor(0, 0, 0)):
    x, y, w, h = bullet_box

    # Positions
    box_width = Inches(x)
    box_height = Inches(y)
    start_left = Inches(w)
    start_top = Inches(h)
    x_spacing = Inches(0.3)
    y_spacing = Inches(0.5)

    for idx, feature in enumerate(items):
        row = idx // 2
        col = idx % 2
        left = start_left + col * (box_width + x_spacing)
        top = start_top + row * (box_height + y_spacing)

        # Shape background
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_width, box_height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 235)  # soft beige
        shape.line.color.rgb = RGBColor(230, 230, 230)

        # Text content
        text_frame = shape.text_frame
        text_frame.clear()
        
        # title_para = text_frame.paragraphs[0]
        # title_para.text = feature['title']
        # title_para.font.bold = True
        # title_para.font.size = Pt(16)
        # title_para.font.color.rgb = RGBColor(67, 94, 71)
        # title_para.alignment = PP_ALIGN.LEFT

        desc_para = text_frame.add_paragraph()
        desc_para.text = feature
        desc_para.font.size = Pt(12)
        desc_para.font.color.rgb = RGBColor(50, 50, 50)
 

def add_paragraph(slide, text, inches, font_size=22, color=RGBColor(0, 0, 0), bold=False):
    para_box = slide.shapes.add_textbox(*map(Inches, inches))
    para_frame = para_box.text_frame
    para_frame.clear()
    para_frame.word_wrap = True
    para_frame.vertical_anchor = MSO_ANCHOR.TOP

    p = para_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold

    para_frame.margin_left = Inches(0.1)
    para_frame.margin_right = Inches(0.1)



def add_image_placeholder(slide, image_prompt, inches):
    placeholder = slide.shapes.add_textbox(*map(Inches, inches))
    tf = placeholder.text_frame
    tf.text = f"[Image: {image_prompt}]"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER



def shapes_adding(slide, x, y, w, h, text, font_size=20, color=RGBColor(0, 0, 0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUND_2_DIAG_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 230, 250)  # light lavender
    shape.line.color.rgb = RGBColor(180, 180, 180)

    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = False
    text_frame.word_wrap = True


def add_numbered_shapes(slide, items, base_inches, font_size=20, color=RGBColor(0, 0, 0)):
    x, y, w, h = base_inches
    line_spacing = h + 0.1  # vertical spacing between shapes
    for idx, text in enumerate(items):
        full_text = text
        shapes_adding(slide, x, y + idx * line_spacing, w, h, full_text, font_size=font_size, color=color)
