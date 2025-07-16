from django.views.decorators.csrf import csrf_exempt
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
from rest_framework.decorators import api_view, permission_classes
from rest_framework.permissions import AllowAny
from rest_framework.response import Response
from rest_framework import status
from pptx import Presentation
from .serializers import PresentationGenerationSerializer
from .Utils.Ai_text_engine import genrate_gemni,generate_openai,generate_haiku
from django.http import HttpResponse, JsonResponse, HttpResponseRedirect
from django.shortcuts import render, redirect
from django.urls import reverse
from pptx import Presentation
from django.core.files.base import ContentFile
import io
import json
from User.models import GeneratedPresentation
from .Utils.Pptgen import create_slide
import uuid
import re
import io
import os
from datetime import datetime



# In-memory stores (for demo)
user_preferences, user_prompts, user_attachments, analytics_data = {}, {}, {},{}

# # ========================== BASIC USER DATA ENDPOINTS ==========================

@api_view(['GET'])
@permission_classes([AllowAny])
def get_user_preferences(request):
    user_id = request.GET.get('userId', 'default_user')
    try:
        return Response({
            'lastSelectedOptions': user_preferences.get(user_id, {
                'Model': 'Gemini',
                'Slides': '5',
                'Font': 'Arial',
                'Style': 'Professional',
                'Content': 'Medium'
            }),
            'lastPrompt': user_prompts.get(user_id, ''),
            'savedAttachments': user_attachments.get(user_id, [])
        })
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['POST'])
@permission_classes([AllowAny])
def save_user_preferences(request):
    try:
        user_id = request.data.get('userId', 'default_user')
        user_preferences[user_id] = request.data.get('selectedOptions', {})
        return Response({'success': True, 'message': 'Preferences saved'})
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['POST'])
@permission_classes([AllowAny])
def auto_save_prompt(request):
    try:
        user_id = request.data.get('userId', 'default_user')
        user_prompts[user_id] = request.data.get('prompt', '')
        return Response({'success': True, 'message': 'Prompt saved'})
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# ========================== ANALYTICS TRACKERS ==========================

def log_event(event_type, data):
    data.update({
        'type': event_type,
        'timestamp': data.get('timestamp', datetime.now().isoformat()),
        'id': str(uuid.uuid4())
    })
    analytics_data.append(data)  


@api_view(['POST'])
@permission_classes([AllowAny])
def log_analytics_error(request):
    try:
        log_event('error', {
            'error': request.data.get('error', ''),
            'userId': request.data.get('userId', '')
        })
        return Response({'success': True, 'message': 'Error logged'})
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['POST'])
@permission_classes([AllowAny])
def track_option_usage(request):
    try:
        log_event('option_usage', {
            'optionType': request.data.get('optionType', ''),
            'selectedValue': request.data.get('selectedValue', ''),
            'userId': request.data.get('userId', ''),
            'sessionId': request.data.get('sessionId', '')
        })
        return Response({'success': True, 'message': 'Option usage tracked'})
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['POST'])
@permission_classes([AllowAny])
def track_feature_usage(request):
    try:
        log_event('feature_usage', {
            'feature': request.data.get('feature', ''),
            'userId': request.data.get('userId', '')
        })
        return Response({'success': True, 'message': 'Feature usage tracked'})
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['POST'])
@permission_classes([AllowAny])
def track_navigation(request):
    try:
        log_event('navigation', {
            'action': request.data.get('action', '')
        })
        return Response({'success': True, 'message': 'Navigation tracked'})
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['GET'])
@permission_classes([AllowAny])
def get_analytics(request):
    try:
        return Response({
            'analytics': analytics_data,
            'total_events': len(analytics_data),
            'users': len(user_preferences),
            'attachments': sum(len(v) for v in user_attachments.values())
        })
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# ==================================================== FILE UPLOADS =============================================================

@api_view(['POST'])
@permission_classes([AllowAny])
def upload_file(request):
    try:
        if 'file' not in request.FILES:
            return Response({'error': 'No file provided'}, status=status.HTTP_400_BAD_REQUEST)

        file = request.FILES['file']
        user_id = request.POST.get('userId', 'default_user')
        upload_type = request.POST.get('uploadType', 'general_file')

        filename = f"{uuid.uuid4()}.{file.name.split('.')[-1]}"
        path = default_storage.save(f"uploads/{filename}", ContentFile(file.read()))
        file_url = f"/media/{path}"

        file_info = {
            'id': str(uuid.uuid4()),
            'name': file.name,
            'url': file_url,
            'size': file.size,
            'mimeType': file.content_type,
            'uploadType': upload_type,
            'userId': user_id,
            'timestamp': datetime.now().isoformat()
        }

        user_attachments.setdefault(user_id, []).append(file_info)
        return Response({'success': True, 'fileId': file_info['id'], 'fileUrl': file_url})
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['DELETE'])
@permission_classes([AllowAny])
def remove_attachment(request, attachment_id):
    try:
        for user_id in user_attachments:
            user_attachments[user_id] = [f for f in user_attachments[user_id] if f['id'] != attachment_id]
        return Response({'success': True, 'message': 'Attachment removed'})
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    


# ========================== AI Dispatcher Content writing==========================

def generate_ai_content(model_name, prompt, slides,style,level):
    model_name = model_name.lower()

    if model_name == 'openai':
        return generate_openai(prompt, slides)
    elif model_name == 'claudehaiku':
        return generate_haiku(prompt, slides)
    elif model_name == 'gemini':
        return genrate_gemni(prompt, slides,style,level)
    else:
        return "Invalid model selected."


from django.http import JsonResponse
import json, re
from .Utils.Text_cleaning import parse_ai_raw_text

# ========================== Presentation API ==========================

@api_view(['POST'])
@csrf_exempt
@permission_classes([AllowAny])
def generate_presentation(request):
    serializer = PresentationGenerationSerializer(data=request.data)

    if not serializer.is_valid():
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

    try:
        data = serializer.validated_data
        prompt = data['prompt']
        selected_options = data.get('selectedOptions', {})

        model_name = selected_options.get('Model', 'Gemini')
        slides = selected_options.get('Slides', '5')
        # Font = selected_options.get('Font' , 'Arial')
        style = selected_options.get('Style', 'Professional')
        level = selected_options.get('Content', 'Medium')

        # This field is kept but not used in generation
        attachments = data.get('attachments', [])
        print(prompt)
        print(user_attachments)
        print(user_preferences)
        print(user_prompts)
        ai_raw_text = generate_ai_content(model_name, prompt, slides,style,level)

        cleaned = parse_ai_raw_text(ai_raw_text)

        if ai_raw_text.startswith("[ERROR]"):
                return Response({"error": ai_raw_text}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
        
        cleaned = re.sub(r'^```json\n|```$', '', ai_raw_text.strip(), flags=re.MULTILINE)
        
        # cont = parse_ai_content_to_slides(ai_raw_text)
        # print(slides)
        # return Response({"slides": slides}, status=status.HTTP_200_OK)
    
        slides = json.loads(cleaned)
        return JsonResponse({
            'success': True,
            'presentation': {
                'id': str(uuid.uuid4()),
                'title': f"Presentation - {prompt[:30]}...",
                "slides": slides,
                'metadata': selected_options,
            }
        })
        


    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


    

# =============================  AI text Genration Models ===========================================


def generate_fallback_content(prompt, slides):
    return "\n\n".join(
        f"SLIDE {i}: {prompt} Part {i}\n- Point 1\n- Point 2\n- Point 3"
        for i in range(1, int(slides) + 1)
    )

def parse_ai_content_to_slides(content):
    slides = []
    for block in content.split('\n\n'):
        lines = block.strip().split('\n')
        if not lines: continue
        title = lines[0].split(':', 1)[-1].strip() if 'SLIDE' in lines[0] else lines[0]
        bullets = [line.lstrip('-• ').strip() for line in lines[1:] if line.strip()]
        slides.append({
            'id': str(uuid.uuid4()),
            'title': title,
            'content': bullets,
            'slide_number': len(slides) + 1
        })
    return slides




# # # ====================================== creating Presentation =============================================


def generate_ppt(request):
    if request.method == 'POST':
        if request.user.is_authenticated:
            raw_data = request.POST.get('data')  # JSON string
            topic_name = request.POST.get('topic_name')
            theme_name = request.POST.get('theme_name')

            if not raw_data or not topic_name:
                return HttpResponse("Missing data or topic name", status=400)

            try:
                data = json.loads(raw_data)
            except json.JSONDecodeError:
                return HttpResponse("Invalid JSON data", status=400)

            prs = Presentation("Light_pink.pptx")  # your base template

            for slide_data in data:
                create_slide(prs, slide_data)

            # Save the presentation to a file
            ppt_io = io.BytesIO()
            prs.save(ppt_io)
            ppt_io.seek(0)

            filename = f"{topic_name}_presentation.pptx"
            file_content = ContentFile(ppt_io.read(), name=filename)

            # Save to database
            GeneratedPresentation.objects.create(
                user=request.user,
                file=file_content,
                topic_name=topic_name
            )

            # Return downloadable response
            ppt_io.seek(0)
            response = HttpResponse(
                ppt_io.getvalue(),
                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            response['Content-Disposition'] = f'attachment; filename="{topic_name}_presentation.pptx"'
            return response

        else:
            if request.headers.get('x-requested-with') == 'XMLHttpRequest':
                return JsonResponse({'error': 'User not authenticated'}, status=401)
            else:
                # Redirect browser to login page
                return redirect(f"{reverse('login')}?next={request.path}")

    # For GET requests
    return render(request, 'generate_ppt.html')




# @csrf_exempt
# def generate_ppt(request):
#     if request.method == 'POST':
#         topic = request.POST.get('topic', 'AI')
#         api = request.POST.get('api', 'gemini').lower()
#         # content = genrate_ai_content(topic, api, '5', 'Professional', 'Medium')
#         if not content:
#             return HttpResponse("AI content generation failed.")

#         slides = parse_ai_content_to_slides(content)
#         prs = Presentation()
#         for slide in slides:
#             layout = prs.slide_layouts[1]
#             s = prs.slides.add_slide(layout)
#             s.shapes.title.text = slide['title']
#             tf = s.placeholders[1].text_frame
#             tf.clear()
#             for point in slide['content']:
#                 tf.add_paragraph().text = point

#         ppt_file = io.BytesIO()
#         prs.save(ppt_file)
#         ppt_file.seek(0)

#         response = HttpResponse(ppt_file.getvalue(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
#         response['Content-Disposition'] = f'attachment; filename="{topic}_presentation.pptx"'
#         return response
#     return render(request, 'generate_ppt.html')